from pathlib import Path
from pptx import Presentation
from src.parsers.base import DocumentParser, ParsedDocument


class PptxParser(DocumentParser):
    def parse(self, path: Path) -> ParsedDocument:
        prs = Presentation(str(path))
        slides_text = []
        for i, slide in enumerate(prs.slides):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    texts.append(shape.text_frame.text)
            if texts:
                slides_text.append(f"## Slide {i + 1}\n" + "\n".join(texts))
        return ParsedDocument(
            content="\n\n".join(slides_text),
            metadata={"slides": len(slides_text)},
            source_path=str(path),
            file_type="pptx",
        )
